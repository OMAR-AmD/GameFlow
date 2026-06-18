from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialiser la présentation
prs = Presentation()

# Fonctions utilitaires
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Titre
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    return slide

def add_bullet_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Titre et Contenu
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if ":" in point and not point.startswith("  "):
            pass
        elif point.startswith("-"):
            p.level = 1
            
    return slide

# Slide 1 : Titre
add_title_slide(
    prs, 
    "GameFlow AI", 
    "Behavior-Driven Game Recommendation System\n\nPréparé par : Akby Anass & Omar Amdouni\nEHTP - MIG (2025-2026)"
)

# Slide 2 : Le Problème
add_bullet_slide(
    prs,
    "1. Contexte & Problématique",
    [
        "Le défi de la plateforme Steam : plus de 50 000 jeux vidéo disponibles.",
        "Limites des systèmes classiques :",
        "- Se basent souvent uniquement sur les tags (RPG, Action).",
        "- Incapables de différencier deux joueurs de RPG (un complétionniste vs un explorateur).",
        "Notre objectif :",
        "- Créer un moteur de recommandation basé sur le comportement réel (temps de jeu)."
    ]
)

# Slide 3 : Feature Engineering
add_bullet_slide(
    prs,
    "2. Données & Feature Engineering",
    [
        "Les données : Dataset Steam-200k (Interactions) & Steam Metadata (Genres).",
        "Le temps de jeu comme proxy comportemental fort.",
        "Création de métriques sur-mesure :",
        "- Session Intensity (Temps de jeu / Médiane)",
        "- Completion Ratio",
        "- Competitive Index",
        "- Abandonment Rate (Taux d'abandon avant 1 heure)"
    ]
)

# Slide 4 : Apprentissage Non-Supervisé
add_bullet_slide(
    prs,
    "3. Le Profilage : Clustering K-Means",
    [
        "Application de PCA et K-Means sur les métriques comportementales.",
        "Extraction de 5 Personas clés :",
        "1. Compétiteur Hardcore (47% - Multijoueur intense)",
        "2. Collectionneur Versatile (19.5% - Touche à tout)",
        "3. Zappeur Curieux (17% - Abandon rapide)",
        "4. Marathonien Passionné (8.5% - Sessions très longues)",
        "5. Explorateur Narratif (8% - Focus sur l'histoire)"
    ]
)

# Slide 5 : Le Modèle Hybride
add_bullet_slide(
    prs,
    "4. Le Moteur de Recommandation Hybride",
    [
        "Filtrage Collaboratif (SVD) :",
        "- Factorisation de la matrice Joueur-Jeu (TruncatedSVD).",
        "- Identification des goûts latents.",
        "Filtrage Basé sur le Contenu (TF-IDF) :",
        "- Similarité cosinus sur les genres des jeux.",
        "La fusion (Hybride) :",
        "- Score = α * SVD + (1-α) * TF-IDF",
        "- Permet de résoudre le problème du démarrage à froid (Cold Start)."
    ]
)

# Slide 6 : L'Application Web
add_bullet_slide(
    prs,
    "5. Application Full-Stack & Explainable AI",
    [
        "Architecture moderne : Base SQLite + API FastAPI + Frontend React/Vite.",
        "Solveur de 'Cold Start' innovant :",
        "- Déduction de la Persona via un questionnaire initial.",
        "L'IA Explicable (Explainable AI) :",
        "- Justification claire de chaque recommandation.",
        "- Exemple : 'Score Collaboratif 85% - Score Contenu 15%'.",
        "Dashboard Analytique Administrateur (Recharts)."
    ]
)

# Slide 7 : Conclusion
add_bullet_slide(
    prs,
    "6. Conclusion",
    [
        "Le comportement (temps de jeu) est plus riche que le simple genre.",
        "Dépassement du cahier des charges initial :",
        "- Produit final déployable (FastAPI + React).",
        "- Tableaux de bord télémétriques.",
        "Perspectives :",
        "- Deep Learning (RNN) pour l'ordre chronologique des jeux.",
        "- Déploiement Cloud via Docker."
    ]
)

# Sauvegarder la présentation
output_filename = "GameFlow_Presentation.pptx"
prs.save(output_filename)
print(f"Présentation générée avec succès : {output_filename}")
